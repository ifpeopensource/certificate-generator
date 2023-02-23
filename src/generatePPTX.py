from pathlib import Path
from typing import Dict, List
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def handleAlignOption(align: str):
    if align == "left":
        return PP_ALIGN.LEFT
    elif align == "center":
        return PP_ALIGN.CENTER
    elif align == "right":
        return PP_ALIGN.RIGHT
    elif align == "justify":
        return PP_ALIGN.JUSTIFY
    else:
        return None

def generatePPTX(model: str, fields: List[str], data: List[str], options, output_dir: str) -> Path:
    align = handleAlignOption(options['align'])
    prs = Presentation(model)
    slide = prs.slides[0]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        frame = shape.text_frame

        for field_index, field in enumerate(fields):
            field_placeholder = f"{{{{{field}}}}}"
            if field_placeholder in frame.text:
                frame.text = frame.text.replace(field_placeholder, data[field_index])

                for paragraph in frame.paragraphs:
                    paragraph.alignment = align
                    paragraph.font.size = Pt(options['font_size'])
                    paragraph.font.color.rgb = RGBColor.from_string(
                        options['color'])

    file_path = Path(output_dir).joinpath(f"pptx/{data[0]}.pptx")
    prs.core_properties.title = f"Certificate for {data[0]}"
    prs.core_properties.author = "IFPE Open Source"
    prs.core_properties.comments = "Certificate generated using IFPE Open Source Certificate Generator\nCertificado gerado usando o IFPE Open Source Certificate Generator\nhttps://github.com/ifpeopensource/certificate-generator"
    prs.save(file_path)
    return file_path
